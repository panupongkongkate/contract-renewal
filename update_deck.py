from pathlib import Path
from shutil import copyfile

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


SRC = Path(r"C:\Users\Ment\Desktop\Presentation\contract-renewal-panupong\Panupong-Contract-Renewal-v3.pptx")
DST = Path(r"C:\Users\Ment\Desktop\Presentation\contract-renewal-panupong\Panupong-Contract-Renewal-v4.pptx")


def fit_picture(shape, left, top, max_w, max_h):
    ratio = shape.width / shape.height
    box_ratio = max_w / max_h
    if ratio > box_ratio:
        new_w = max_w
        new_h = int(max_w / ratio)
    else:
        new_h = max_h
        new_w = int(max_h * ratio)
    shape.width = int(new_w)
    shape.height = int(new_h)
    shape.left = int(left + (max_w - new_w) / 2)
    shape.top = int(top + (max_h - new_h) / 2)


def _clear_bullets(paragraph):
    p_pr = paragraph._p.get_or_add_pPr()
    for tag in ("buClrTx", "buSzTx", "buFontTx", "buNone", "buChar", "buAutoNum"):
        node = p_pr.find(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}")
        if node is not None:
            p_pr.remove(node)
    p_pr.append(OxmlElement("a:buNone"))


def set_single_paragraph(shape, text, size_pt, bold=False, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _clear_bullets(p)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(42, 42, 42)
    p.alignment = align


def set_multi_paragraph(shape, lines, size_pt, bold=False, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _clear_bullets(p)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(42, 42, 42)
        p.alignment = align


def make_text_transparent(shape):
    if hasattr(shape, "fill"):
        shape.fill.background()
    if hasattr(shape, "line"):
        shape.line.fill.background()


def set_title(shape, text, top=1.22):
    shape.left = Inches(0.95)
    shape.top = Inches(top)
    shape.width = Inches(11.0)
    shape.height = Inches(0.42)
    set_single_paragraph(shape, text, 22, bold=True, align=PP_ALIGN.LEFT)


def set_caption(shape, text, left, top, width, size=12):
    shape.left = left
    shape.top = top
    shape.width = width
    shape.height = Inches(0.24)
    set_single_paragraph(shape, text, size, align=PP_ALIGN.CENTER)


copyfile(SRC, DST)
prs = Presentation(str(DST))

# Slide 6
slide = prs.slides[5]
set_title(slide.shapes[5], "โครงการ newKronos (ลงทะเบียนล่วงหน้า)")
fit_picture(slide.shapes[6], Inches(5.95), Inches(1.18), Inches(6.55), Inches(6.08))

# Slide 7
slide = prs.slides[6]
set_title(slide.shapes[5], "โครงการ newKronos (ลงทะเบียนเข้าอาคาร)")
fit_picture(slide.shapes[6], Inches(0.88), Inches(1.82), Inches(11.60), Inches(5.38))

# Slide 8
slide = prs.slides[7]
set_title(slide.shapes[5], "โครงการ newKronos (Report)", top=1.00)
fit_picture(slide.shapes[6], Inches(0.90), Inches(1.72), Inches(5.40), Inches(2.55))
set_caption(slide.shapes[7], "รายงานเข้า-ออกประจำวัน (เว็บ)", Inches(1.00), Inches(4.30), Inches(5.20))
fit_picture(slide.shapes[8], Inches(6.40), Inches(1.72), Inches(5.40), Inches(2.55))
set_caption(slide.shapes[9], "รายงานสรุปเข้า-ออก (เว็บ)", Inches(6.50), Inches(4.30), Inches(5.20))
fit_picture(slide.shapes[10], Inches(0.90), Inches(4.70), Inches(5.40), Inches(2.05))
set_caption(slide.shapes[11], "รายงานเข้า-ออกประจำวัน (Excel)", Inches(1.00), Inches(6.83), Inches(5.20))
fit_picture(slide.shapes[12], Inches(6.40), Inches(4.70), Inches(5.40), Inches(2.05))
set_caption(slide.shapes[13], "รายงานสรุปเข้า-ออก (Excel)", Inches(6.50), Inches(6.83), Inches(5.20))

# Slide 9
slide = prs.slides[8]
set_title(slide.shapes[5], "โครงการ newKronos (Master Data)")
fit_picture(slide.shapes[6], Inches(0.90), Inches(1.80), Inches(5.45), Inches(2.55))
set_caption(slide.shapes[7], "เมนู ผู้ติดต่อ", Inches(1.00), Inches(4.38), Inches(5.20))
fit_picture(slide.shapes[8], Inches(6.35), Inches(1.80), Inches(5.45), Inches(2.55))
set_caption(slide.shapes[9], "เมนู อาคาร/ตึก", Inches(6.45), Inches(4.38), Inches(5.20))
fit_picture(slide.shapes[10], Inches(0.90), Inches(4.72), Inches(5.45), Inches(2.00))
set_caption(slide.shapes[11], "เมนู ชั้น", Inches(1.00), Inches(6.83), Inches(5.20))
fit_picture(slide.shapes[12], Inches(6.35), Inches(4.72), Inches(5.45), Inches(2.00))
set_caption(slide.shapes[13], "เมนู บริษัท", Inches(6.45), Inches(6.83), Inches(5.20))

# Slide 10
slide = prs.slides[9]
set_title(slide.shapes[5], "โครงการ newKronos (Master Data)")
fit_picture(slide.shapes[6], Inches(0.90), Inches(1.80), Inches(5.45), Inches(2.55))
set_caption(slide.shapes[7], "เมนู วัตถุประสงค์", Inches(1.00), Inches(4.38), Inches(5.20))
fit_picture(slide.shapes[8], Inches(6.35), Inches(1.80), Inches(5.45), Inches(2.55))
set_caption(slide.shapes[9], "เมนู พนักงาน", Inches(6.45), Inches(4.38), Inches(5.20))
fit_picture(slide.shapes[10], Inches(0.90), Inches(4.72), Inches(5.45), Inches(2.00))
set_caption(slide.shapes[11], "เมนู บัตรแลกเข้า", Inches(1.00), Inches(6.83), Inches(5.20))
fit_picture(slide.shapes[12], Inches(6.35), Inches(4.72), Inches(5.45), Inches(2.00))
set_caption(slide.shapes[13], "เมนู จัดการพื้นที่", Inches(6.45), Inches(6.83), Inches(5.20))

# Slide 11
slide = prs.slides[10]
set_title(slide.shapes[5], "โครงการ DataCenter")
slide.shapes[6].left = Inches(0.95)
slide.shapes[6].top = Inches(1.72)
slide.shapes[6].width = Inches(8.20)
slide.shapes[6].height = Inches(0.28)
set_single_paragraph(slide.shapes[6], "สรุปบทบาทและผลลัพธ์ของงาน BackEnd ในโครงการ DataCenter", 12, align=PP_ALIGN.LEFT)
slide.shapes[7].left = Inches(0.95)
slide.shapes[7].top = Inches(2.30)
slide.shapes[7].width = Inches(4.55)
slide.shapes[7].height = Inches(1.85)
slide.shapes[8].left = Inches(1.18)
slide.shapes[8].top = Inches(2.72)
slide.shapes[8].width = Inches(3.00)
slide.shapes[8].height = Inches(0.26)
make_text_transparent(slide.shapes[8])
set_single_paragraph(slide.shapes[8], "บทบาทที่รับผิดชอบ", 18, bold=True, align=PP_ALIGN.LEFT)
slide.shapes[9].left = Inches(1.18)
slide.shapes[9].top = Inches(3.14)
slide.shapes[9].width = Inches(3.95)
slide.shapes[9].height = Inches(0.70)
make_text_transparent(slide.shapes[9])
set_multi_paragraph(
    slide.shapes[9],
    [
        "ดูแลการพัฒนา BackEnd",
        "และจัดทำ API ที่ระบบต้องการใช้งาน",
    ],
    16,
    align=PP_ALIGN.LEFT,
)
slide.shapes[10].left = Inches(5.70)
slide.shapes[10].top = Inches(2.30)
slide.shapes[10].width = Inches(4.55)
slide.shapes[10].height = Inches(1.85)
slide.shapes[11].left = Inches(5.95)
slide.shapes[11].top = Inches(2.72)
slide.shapes[11].width = Inches(2.70)
slide.shapes[11].height = Inches(0.26)
make_text_transparent(slide.shapes[11])
set_single_paragraph(slide.shapes[11], "งานที่ดำเนินการ", 18, bold=True, align=PP_ALIGN.LEFT)
slide.shapes[12].left = Inches(5.95)
slide.shapes[12].top = Inches(3.14)
slide.shapes[12].width = Inches(3.95)
slide.shapes[12].height = Inches(0.70)
make_text_transparent(slide.shapes[12])
set_multi_paragraph(
    slide.shapes[12],
    [
        "พัฒนา API ครอบคลุมงานหลักของระบบ",
        "และจัดทำกลุ่มรายงาน",
    ],
    16,
    align=PP_ALIGN.LEFT,
)
slide.shapes[13].left = Inches(0.95)
slide.shapes[13].top = Inches(4.55)
slide.shapes[13].width = Inches(10.55)
slide.shapes[13].height = Inches(1.45)
slide.shapes[14].left = Inches(1.20)
slide.shapes[14].top = Inches(4.92)
slide.shapes[14].width = Inches(9.95)
slide.shapes[14].height = Inches(0.88)
make_text_transparent(slide.shapes[14])
set_multi_paragraph(
    slide.shapes[14],
    [
        "รายงานที่พัฒนา: กลุ่มรายงานการเข้าออกอาคาร และกลุ่มรายงานการลงทะเบียนเข้าอาคาร",
        "ผลลัพธ์: รองรับการลงทะเบียนและเข้าอาคารได้รวดเร็วขึ้น",
    ],
    16,
    align=PP_ALIGN.LEFT,
)

# Slide 12
slide = prs.slides[11]
set_title(slide.shapes[5], "โครงการ DataCenter (ลงทะเบียนล่วงหน้า)")
fit_picture(slide.shapes[6], Inches(0.85), Inches(2.18), Inches(5.65), Inches(3.30))
set_caption(slide.shapes[7], "ลงทะเบียนล่วงหน้า Step 1", Inches(1.10), Inches(5.78), Inches(5.15), 13)
fit_picture(slide.shapes[8], Inches(6.60), Inches(2.18), Inches(5.65), Inches(3.30))
set_caption(slide.shapes[9], "ลงทะเบียนล่วงหน้า Step 2", Inches(6.85), Inches(5.78), Inches(5.15), 13)

# Slide 13
slide = prs.slides[12]
set_title(slide.shapes[5], "โครงการ DataCenter (ลงทะเบียนเข้าอาคาร)")
set_single_paragraph(slide.shapes[6], "", 12, align=PP_ALIGN.LEFT)
fit_picture(slide.shapes[7], Inches(0.95), Inches(1.82), Inches(10.10), Inches(5.45))
slide.shapes[8].left = Inches(11.18)
slide.shapes[8].top = Inches(2.30)
slide.shapes[8].width = Inches(1.20)
slide.shapes[8].height = Inches(2.10)
set_multi_paragraph(slide.shapes[8], ["ลงทะเบียนเข้าอาคาร", "ใช้ PIN code", "ในการเข้าอาคาร"], 18, bold=True, align=PP_ALIGN.LEFT)

# Slide 14
slide = prs.slides[13]
set_title(slide.shapes[4], "โครงการ DataCenter (รายการลงทะเบียนล่วงหน้า)")
fit_picture(slide.shapes[5], Inches(0.95), Inches(1.72), Inches(7.55), Inches(5.55))
slide.shapes[6].text_frame.margin_left = 0
slide.shapes[6].text_frame.margin_right = 0
slide.shapes[6].text_frame.margin_top = Pt(4)
slide.shapes[6].text_frame.margin_bottom = Pt(4)
slide.shapes[6].left = Inches(8.78)
slide.shapes[6].top = Inches(2.18)
slide.shapes[6].width = Inches(3.45)
slide.shapes[6].height = Inches(3.70)
set_multi_paragraph(
    slide.shapes[6],
    [
        "ฟังก์ชันหลัก",
        "- แก้ไขข้อมูล",
        "- คัดลอกใบงาน",
        "- ลบใบงาน",
        "- ค้นหาใบงาน",
        "- อนุมัติใบงาน",
    ],
    20,
    bold=True,
    align=PP_ALIGN.LEFT,
)

# Slide 15
slide = prs.slides[14]
set_title(slide.shapes[5], "โครงการ DataCenter (Report)")
fit_picture(slide.shapes[6], Inches(0.90), Inches(1.72), Inches(5.40), Inches(2.55))
set_caption(slide.shapes[7], "รายงานเข้า-ออกประจำวัน (เว็บ)", Inches(1.00), Inches(4.30), Inches(5.20))
fit_picture(slide.shapes[8], Inches(6.40), Inches(1.72), Inches(5.40), Inches(2.55))
set_caption(slide.shapes[9], "รายงานเข้า-ออก (เว็บ)", Inches(6.50), Inches(4.30), Inches(5.20))
fit_picture(slide.shapes[10], Inches(0.90), Inches(4.70), Inches(5.40), Inches(2.05))
set_caption(slide.shapes[11], "รายงานเข้า-ออกประจำวัน (Excel)", Inches(1.00), Inches(6.83), Inches(5.20))
fit_picture(slide.shapes[12], Inches(6.40), Inches(4.70), Inches(5.40), Inches(2.05))
set_caption(slide.shapes[13], "รายงานเข้า-ออก (Excel)", Inches(6.50), Inches(6.83), Inches(5.20))

# Slide 16
slide = prs.slides[15]
set_title(slide.shapes[5], "โครงการ DataCenter (Report)")
fit_picture(slide.shapes[6], Inches(0.80), Inches(2.00), Inches(5.75), Inches(3.60))
set_caption(slide.shapes[7], "รายการเข้า-ออกพื้นที่ (เว็บ)", Inches(1.10), Inches(5.82), Inches(5.15), 13)
fit_picture(slide.shapes[8], Inches(6.70), Inches(2.00), Inches(5.75), Inches(3.60))
set_caption(slide.shapes[9], "รายการเข้า-ออกพื้นที่ (Excel)", Inches(7.00), Inches(5.82), Inches(5.15), 13)

# Slide 17
slide = prs.slides[16]
set_title(slide.shapes[5], "โครงการ DataCenter (Master Data)")
fit_picture(slide.shapes[6], Inches(0.90), Inches(1.80), Inches(5.45), Inches(2.55))
set_caption(slide.shapes[7], "เมนู วัตถุประสงค์", Inches(1.00), Inches(4.38), Inches(5.20))
fit_picture(slide.shapes[8], Inches(6.35), Inches(1.80), Inches(5.45), Inches(2.55))
set_caption(slide.shapes[9], "เมนู พนักงาน", Inches(6.45), Inches(4.38), Inches(5.20))
fit_picture(slide.shapes[10], Inches(0.90), Inches(4.72), Inches(5.45), Inches(2.00))
set_caption(slide.shapes[11], "เมนู บัตรแลกเข้า", Inches(1.00), Inches(6.83), Inches(5.20))
fit_picture(slide.shapes[13], Inches(6.35), Inches(4.72), Inches(5.45), Inches(2.00))
set_caption(slide.shapes[12], "เมนู จัดการพื้นที่", Inches(6.45), Inches(6.83), Inches(5.20))

# Slide 18
slide = prs.slides[17]
set_title(slide.shapes[5], "โครงการ DataCenter (Master Data)")
fit_picture(slide.shapes[10], Inches(0.85), Inches(1.72), Inches(5.55), Inches(2.65))
set_caption(slide.shapes[6], "เมนู จัดการข้อมูล Rack", Inches(1.00), Inches(4.42), Inches(5.20))
fit_picture(slide.shapes[11], Inches(6.25), Inches(1.72), Inches(5.55), Inches(2.65))
set_caption(slide.shapes[7], "เมนู จัดการข้อมูลห้อง", Inches(6.40), Inches(4.42), Inches(5.20))
fit_picture(slide.shapes[8], Inches(3.55), Inches(4.85), Inches(6.25), Inches(1.95))
set_caption(slide.shapes[9], "เมนู จัดการ permission", Inches(4.10), Inches(6.86), Inches(5.15))

# Slide 22
slide = prs.slides[21]
set_multi_paragraph(slide.shapes[9], ["ตรวจสอบปัญหาลงทะเบียนล่วงหน้าไม่ได้", "จาก ACC หรือข้อมูล company ผิด"], 15, align=PP_ALIGN.LEFT)
set_multi_paragraph(slide.shapes[12], ["ตรวจสอบ Get Company, companyID", "และ token หมดอายุผิดปกติ"], 15, align=PP_ALIGN.LEFT)
set_multi_paragraph(slide.shapes[15], ["ตรวจสอบข้อมูลผู้เข้าอาคาร", "และวิเคราะห์การเชื่อมต่อ ProWatch กับ Access Control"], 15, align=PP_ALIGN.LEFT)
slide.shapes[16].height = Inches(1.20)
set_single_paragraph(slide.shapes[17], "", 14, align=PP_ALIGN.LEFT)
slide.shapes[18].top = Inches(4.76)
body = slide.shapes.add_textbox(Inches(1.22), Inches(5.18), Inches(6.25), Inches(0.72))
body.text_frame.margin_left = 0
body.text_frame.margin_right = 0
body.text_frame.margin_top = 0
body.text_frame.margin_bottom = 0
set_multi_paragraph(
    body,
    [
        "แก้ไขข้อมูลไม่ตรง เช่น ชื่อพนักงาน ผู้อนุมัติ และตรวจสอบใบงานที่ค้นหาไม่พบ",
        "ครอบคลุมทั้งงาน BackEnd และการตรวจสอบ integration หลายระบบ",
    ],
    14,
    align=PP_ALIGN.LEFT,
)

prs.save(str(DST))
